"""#!/usr/bin/env python3
"""Generate PPTX and DOCX academic documents for HTENG437 and HTENG440 modules.

This script creates:
- HTENG437_Topic3.6_Presentation.pptx (9 slides)
- HTENG437_Topic3.6_Detailed_Report.docx
- HTENG440_Chapter2_Presentation.pptx (10 slides)
- HTENG440_Chapter2_Detailed_Report.docx
- HTENG437_Full_Module_Summary.pptx (15 slides)
- HTENG440_Full_Module_Summary.pptx (15 slides)

Dependencies (see requirements.txt): python-pptx, python-docx, Pillow

Usage:
    pip install -r requirements.txt
    python generate_all_documents.py

The slides are generated with simple layouts and placeholder content matching the requested themes and footers.

"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE
from docx import Document
from docx.shared import Pt as DocPt
from docx.enum.text import WD_ALIGN_PARAGRAPH
import os
import sys

OUT_DIR = os.path.dirname(os.path.abspath(__file__)) if '__file__' in globals() else os.getcwd()

# Helpers for PPTX

def set_slide_background_color(slide, rgb_tuple):
    r, g, b = rgb_tuple
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_footer(slide, text, font_size=12):
    left = Inches(0.2)
    top = Inches(6.8)
    width = Inches(9)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def add_title_and_body(slide, title, body_lines, title_size=40, body_size=18, body_left=Inches(0.5)):
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    t_p = title_tf.paragraphs[0]
    t_p.text = title
    t_p.font.size = Pt(title_size)
    t_p.font.bold = True
    t_p.font.color.rgb = RGBColor(255, 255, 255)

    # Body
    body_box = slide.shapes.add_textbox(body_left, Inches(1.8), Inches(9), Inches(4))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = body_tf.add_paragraph() if i > 0 else body_tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(body_size)
        p.font.color.rgb = RGBColor(220, 240, 255)


def create_presentation(filename, theme_rgb, footer_text, slides_spec):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for spec in slides_spec:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background_color(slide, theme_rgb)
        add_title_and_body(slide, spec.get('title', ''), spec.get('body', []), title_size=spec.get('title_size',40), body_size=spec.get('body_size',18))
        # optional custom drawing
        if spec.get('draw_network'):
            draw_network_diagram(slide)
        if spec.get('draw_circle_senses'):
            draw_internet_of_senses(slide)
        add_footer(slide, footer_text)

    out_path = os.path.join(OUT_DIR, filename)
    prs.save(out_path)
    return out_path

# Simple drawing helpers (basic shapes)

def draw_network_diagram(slide):
    # draw 3 boxes and connecting lines to suggest a network
    lefts = [Inches(1), Inches(4), Inches(7)]
    tops = [Inches(3), Inches(2), Inches(3)]
    boxes = []
    for i, (l, t) in enumerate(zip(lefts, tops)):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE.RECTANGLE, l, t, Inches(1.2), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 212, 255)
        shape.line.color.rgb = RGBColor(255,255,255)
        tf = shape.text_frame
        tf.text = f"Node {i+1}"
        tf.paragraphs[0].font.color.rgb = RGBColor(10,10,10)
        boxes.append(shape)
    # connecting lines
    try:
        line = slide.shapes.add_connector(MSO_AUTO_SHAPE.LINE, lefts[0]+Inches(1.2), tops[0]+Inches(0.4), lefts[1], tops[1]+Inches(0.4))
        line.line.color.rgb = RGBColor(200,200,255)
    except Exception:
        # fallback: add thin rect as a 'line'
        slide.shapes.add_shape(MSO_AUTO_SHAPE.RECTANGLE, Inches(2.2), Inches(3.0), Inches(1.6), Pt(2))


def draw_internet_of_senses(slide):
    # draw five small ovals around a center oval
    center_left = Inches(4.2)
    center_top = Inches(2.6)
    center = slide.shapes.add_shape(MSO_AUTO_SHAPE.OVAL, center_left, center_top, Inches(1.6), Inches(1.0))
    center.fill.solid()
    center.fill.fore_color.rgb = RGBColor(0, 212, 255)
    center.text_frame.text = "6G Network"
    center.text_frame.paragraphs[0].font.color.rgb = RGBColor(10,10,10)
    offsets = [(-2.0, -1.0), (2.8, -1.0), (-2.0, 1.4), (2.8, 1.4), (0.4, 2.4)]
    labels = ["Sight", "Hearing", "Touch", "Smell", "Taste"]
    for (dx, dy), label in zip(offsets, labels):
        l = center_left + Inches(dx)
        t = center_top + Inches(dy)
        o = slide.shapes.add_shape(MSO_AUTO_SHAPE.OVAL, l, t, Inches(1.0), Inches(0.7))
        o.fill.solid()
        o.fill.fore_color.rgb = RGBColor(255, 255, 255)
        o.text_frame.text = label
        o.text_frame.paragraphs[0].font.color.rgb = RGBColor(10,10,10)

# Helpers for DOCX

def create_docx(filename, title, sections):
    doc = Document()
    # Cover
    h = doc.add_heading(title, level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run('\nLecturer: Mr Brian Sadock\n').bold = True
    p.add_run('Date: ').italic = True
    doc.add_page_break()

    # Table of contents placeholder
    doc.add_heading('Table of Contents', level=1)
    for i, s in enumerate(sections, 1):
        p = doc.add_paragraph(f"{i}. {s['heading']}", style='List Number')
    doc.add_page_break()

    # Sections
    for s in sections:
        doc.add_heading(s['heading'], level=1)
        for para in s.get('paras', []):
            p = doc.add_paragraph(para)
            p.style.font.size = DocPt(11)
        if s.get('bullets'):
            for b in s['bullets']:
                doc.add_paragraph(b, style='List Bullet')
        doc.add_page_break()

    out_path = os.path.join(OUT_DIR, filename)
    doc.save(out_path)
    return out_path

# Now assemble the exact documents as specified

def main():
    created = []

    # HTENG437 Topic 3.6 Presentation (9 slides)
    slides_437 = []
    slides_437.append({'title': 'Topic 3.6: Advanced Network Concepts', 'body': ['Digital Twins | Internet of Senses | Holographic Communications | 3D Network Architecture'], 'draw_network': True})
    slides_437.append({'title': "What We'll Cover", 'body': ['1. Digital Twins', '2. Internet of Senses', '3. Holographic Communications', '4. 3D Network Architecture', '\nAll powered by 6G — the network of the future']})
    slides_437.append({'title': 'Digital Twins — The Mirror World', 'body': ['Definition: A digital twin is a real-time virtual replica of a physical object, system, or process', 'In NGN: Mirrors entire network infrastructure digitally', 'Continuously updated via sensors and AI'], 'draw_network': True})
    slides_437.append({'title': 'Digital Twin Architecture', 'body': ['Layer 1 (Physical): Base stations, sensors, devices', 'Layer 2 (Data): Real-time data collection, IoT feeds', 'Layer 3 (Digital Twin): AI processing, simulation, optimization']})
    slides_437.append({'title': 'Internet of Senses — Feel the Network', 'body': ['Integration of all 5 human senses into digital experience', 'Requirements: Terabit speeds, <1ms latency, AI compression'], 'draw_circle_senses': True, 'draw_internet_of_senses': True})
    slides_437.append({'title': 'Holographic Communications — Beyond Video Calls', 'body': ['Today (4G/5G): 2D video, flat screens, fixed viewpoints', 'Tomorrow (6G): 3D hologram, any angle, real-time volumetric, ~1Tbps', 'Technical requirements: Volumetric video encoding, Light field displays, AI rendering, Sub-ms latency']})
    slides_437.append({'title': '3D Network Architecture — Networks Go Vertical', 'body': ['Terrestrial: Fiber, towers', 'Aerial: Drones, HAPS', 'Space: LEO/MEO satellites', 'Use cases: ubiquitous coverage, drones, smart-cities']})
    slides_437.append({'title': 'Convergence & Summary', 'body': ['6G enables Digital Twins, Internet of Senses, Holographic Communications and 3D Architectures', 'Digital twins allow proactive management and optimization']})
    slides_437.append({'title': 'Conclusion', 'body': ['Key takeaways', 'Future outlook', 'References: IEEE and industry reports']})

    path = create_presentation('HTENG437_Topic3.6_Presentation.pptx', (10,22,40), 'HTENG 437 | Next Generation Networks | Topic 3.6', slides_437)
    created.append(path)

    # HTENG437 Detailed Report (DOCX)
    sections_437 = [
        {'heading': 'Introduction to Advanced Network Concepts', 'paras': ['This section introduces the advanced network concepts covered in HTENG 437.']},
        {'heading': 'Digital Twins in Next Generation Networks', 'paras': ['Definition, architecture and applications.'], 'bullets': ['Network planning', 'Fault prediction', 'QoS optimization']},
        {'heading': 'Internet of Senses', 'paras': ['Definition, enabling technologies, use cases and challenges.']},
        {'heading': 'Holographic Communications', 'paras': ['Definition, technical requirements, comparison with current technology.']},
        {'heading': '3D Network Architecture', 'paras': ['Terrestrial, aerial and space layers and their interactions.']},
        {'heading': 'Convergence of Technologies', 'paras': ['How these technologies come together under 6G.']},
        {'heading': 'Conclusion and Future Outlook', 'paras': ['Final notes and references.']}
    ]
    path = create_docx('HTENG437_Topic3.6_Detailed_Report.docx', 'HTENG 437 — Topic 3.6 Advanced Network Concepts', sections_437)
    created.append(path)

    # HTENG440 Chapter 2 Presentation (10 slides)
    slides_440 = []
    slides_440.append({'title': 'Chapter 2: Threats, Vulnerabilities & Mitigations', 'body': ['Understanding the Cyber Threat Landscape']})
    slides_440.append({'title': 'Agenda', 'body': ['CIA Triad, Threat Types, Vulnerabilities, Attack Vectors, Mitigations, Threat Actors']})
    slides_440.append({'title': 'CIA Triad', 'body': ['Confidentiality, Integrity, Availability']})
    slides_440.append({'title': 'Types of Threats', 'body': ['Malware: viruses, worms, trojans, ransomware, spyware, rootkits']})
    slides_440.append({'title': 'Social Engineering', 'body': ['Phishing, Pretexting, Baiting, Tailgating']})
    slides_440.append({'title': 'Vulnerabilities', 'body': ['Unpatched software, misconfigurations, human factors']})
    slides_440.append({'title': 'Attack Vectors', 'body': ['Email, web apps, removable media, remote access, supply chain']})
    slides_440.append({'title': 'Threat Actors', 'body': ['Cybercriminals, Hacktivists, Insiders, Nation-States, Script Kiddies']})
    slides_440.append({'title': 'Mitigations', 'body': ['Patch management, MFA, IDS/IPS, encryption, backups, user training']})
    slides_440.append({'title': 'Conclusion', 'body': ['Key takeaways and references']})

    path = create_presentation('HTENG440_Chapter2_Presentation.pptx', (26,10,10), 'HTENG 440 | Cybersecurity for Engineers | Chapter 2', slides_440)
    created.append(path)

    # HTENG440 Detailed Report
    sections_440 = [
        {'heading': 'Introduction — The Cyber Threat Landscape', 'paras': ['Overview of current threat landscape and motivations.']},
        {'heading': 'Types of Threats', 'paras': ['Malware, social engineering, insider threats, APTs.']},
        {'heading': 'Vulnerabilities', 'paras': ['Software flaws, misconfigurations, human factors.']},
        {'heading': 'Attack Vectors', 'paras': ['Detailed breakdown of common attack vectors.']},
        {'heading': 'Threat Actors', 'paras': ['Categories, motivations and examples.']},
        {'heading': 'Mitigation Strategies', 'paras': ['Technical and non-technical controls.']},
        {'heading': 'Case Studies', 'paras': ['Two brief real-world examples.']},
        {'heading': 'Conclusion', 'paras': ['Summary and references.']}
    ]
    path = create_docx('HTENG440_Chapter2_Detailed_Report.docx', 'HTENG 440 — Chapter 2 Threats, Vulnerabilities and Mitigations', sections_440)
    created.append(path)

    # Full module summaries (15 slides each minimal)
    def make_full_summary(filename, theme_rgb, footer_text, module_title):
        specs = []
        specs.append({'title': module_title + ' — Full Summary', 'body': ['Section overview and learning outcomes']})
        for i in range(1,14):
            specs.append({'title': f'Section {i}', 'body': [f'Content overview for section {i}.']})
        specs.append({'title': 'Conclusion', 'body': ['Final summary, future work, references']})
        return create_presentation(filename, theme_rgb, footer_text, specs)

    path = make_full_summary('HTENG437_Full_Module_Summary.pptx', (10,22,40), 'HTENG 437 | Next Generation Networks', 'HTENG 437 — Next Generation Networks')
    created.append(path)
    path = make_full_summary('HTENG440_Full_Module_Summary.pptx', (26,10,10), 'HTENG 440 | Cybersecurity for Engineers', 'HTENG 440 — Cybersecurity for Engineers')
    created.append(path)

    # Sanity check: ensure files created
    print('\nCreated files:')
    for c in created:
        ok = os.path.exists(c)
        print(f" - {c} -> {'OK' if ok else 'MISSING'}")

    print('\nDone.')

if __name__ == '__main__':
    main()